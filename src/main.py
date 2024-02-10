from pptx import Presentation

# Load the presentation
prs = Presentation("input.pptx")

# Example modification: Add a slide
slide_layout = prs.slide_layouts[0]  # Choosing a slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Added by Pyodide"

# Save the modified presentation
prs.save("output.pptx")
