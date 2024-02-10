from pptx import Presentation
from io import BytesIO

# Load the presentation
# pptx_io = BytesIO()
# pptx_io.write(bytes)
# pptx_io.seek(0)
prs = Presentation()

# Example modification: Add a slide
slide_layout = prs.slide_layouts[0]  # Choosing a slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Added by Pyodide"

# Save the modified presentation
output_io = BytesIO()
prs.save(output_io)
output_io.seek(0)
output_io.read()
